# utils.py
from pdf2docx import parse
from docx2pdf import convert
from PIL import Image
import PyPDF2
from pptx import Presentation
from reportlab.pdfgen import canvas
import io
import os


#document convertor

class FileConverter:
    @staticmethod
    def pdf_to_docx(input_path, output_path):
        parse(input_path, output_path)
        
    @staticmethod
    def docx_to_pdf(input_path, output_path):
        convert(input_path, output_path)
    
    @staticmethod
    def image_convert(input_path, output_path, format):
        image = Image.open(input_path)
        # Convert RGBA to RGB if needed
        if image.mode in ('RGBA', 'LA'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        image.save(output_path, format)
    
    @staticmethod
    def images_to_pdf(input_path, output_path):
        image = Image.open(input_path)
        if image.mode in ('RGBA', 'LA'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        image.save(output_path, 'PDF', resolution=100.0)
        
        
        
#document editor      
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import io
    
        
class DocumentEditor:
    @staticmethod
    def edit_pdf(file_path, pages_to_remove=None, pages_to_add=None, insert_positions=None):
        pdf_reader = PyPDF2.PdfReader(file_path)
        pdf_writer = PyPDF2.PdfWriter()
        
        # Convert to 0-based indexing
        current_pages = list(range(len(pdf_reader.pages)))
        if pages_to_remove:
            for page_num in sorted(pages_to_remove, reverse=True):
                if 0 <= page_num < len(current_pages):
                    current_pages.pop(page_num)

        # Add existing pages
        for page_num in current_pages:
            pdf_writer.add_page(pdf_reader.pages[page_num])

        # Add new pages at specified positions
        if pages_to_add and insert_positions:
            for page, position in zip(pages_to_add, insert_positions):
                if 0 <= position <= len(pdf_writer.pages):
                    pdf_writer.insert_page(page, position)

        output = io.BytesIO()
        pdf_writer.write(output)
        output.seek(0)
        return output

    @staticmethod
    def edit_docx(file_path, pages_to_remove=None, pages_to_add=None, insert_positions=None):
        doc = DocxDocument(file_path)
        
        # Remove paragraphs (approximating pages)
        if pages_to_remove:
            # Calculate approximate paragraphs per page (assuming 40 paragraphs per page)
            paras_per_page = 40
            for page_num in sorted(pages_to_remove, reverse=True):
                start_para = page_num * paras_per_page
                end_para = start_para + paras_per_page
                for i in range(min(end_para, len(doc.paragraphs)) - 1, 
                             start_para - 1, -1):
                    if i < len(doc.paragraphs):
                        p = doc.paragraphs[i]._element
                        p.getparent().remove(p)

        # Add new pages (sections)
        if pages_to_add and insert_positions:
            for content, position in zip(pages_to_add, insert_positions):
                sec = doc.add_section()
                para = sec.add_paragraph(content)

        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    @staticmethod
    def edit_pptx(file_path, slides_to_remove=None, slides_to_add=None, insert_positions=None):
        prs = Presentation(file_path)
        
        # Remove slides
        if slides_to_remove:
            for slide_num in sorted(slides_to_remove, reverse=True):
                if 0 <= slide_num < len(prs.slides):
                    xml_slides = prs.slides._sldIdLst  
                    slides = list(xml_slides)
                    xml_slides.remove(slides[slide_num])

        # Add new slides
        if slides_to_add and insert_positions:
            for content, position in zip(slides_to_add, insert_positions):
                if 0 <= position <= len(prs.slides):
                    blank_slide_layout = prs.slide_layouts[6]  # blank layout
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.title.text = content

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output
